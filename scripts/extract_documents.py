#!/usr/bin/env python3
"""
Document Extraction Pipeline for RAG System
Extracts text from PDF, PPTX, and DOCX files and chunks them for embedding generation.
"""

import os
import json
import hashlib
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass, asdict
import re

# Document processing libraries
try:
    import PyPDF2
    import pdfplumber
except ImportError:
    print("Warning: PDF libraries not installed. Install with: pip install PyPDF2 pdfplumber")

try:
    from pptx import Presentation
except ImportError:
    print("Warning: python-pptx not installed. Install with: pip install python-pptx")

try:
    from docx import Document
except ImportError:
    print("Warning: python-docx not installed. Install with: pip install python-docx")

# Text processing
import tiktoken

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('document_extraction.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)


@dataclass
class ChunkMetadata:
    """Metadata for a text chunk"""
    filename: str
    doc_type: str
    page: Optional[int] = None
    slide: Optional[int] = None
    section: Optional[str] = None
    total_pages: Optional[int] = None
    source_path: Optional[str] = None


@dataclass
class TextChunk:
    """Represents a chunk of text with metadata"""
    id: str
    text: str
    metadata: ChunkMetadata
    token_count: int

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization"""
        return {
            'id': self.id,
            'text': self.text,
            'metadata': asdict(self.metadata),
            'token_count': self.token_count
        }


class DocumentProcessor:
    """Main document processing class"""

    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        """
        Initialize the document processor

        Args:
            chunk_size: Maximum tokens per chunk
            chunk_overlap: Overlap tokens between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.encoding = tiktoken.get_encoding("cl100k_base")  # GPT-3.5/4 encoding
        logger.info(f"Initialized DocumentProcessor (chunk_size={chunk_size}, overlap={chunk_overlap})")

    def count_tokens(self, text: str) -> int:
        """Count tokens in text using tiktoken"""
        return len(self.encoding.encode(text))

    def chunk_text(self, text: str, metadata: ChunkMetadata, doc_id: str) -> List[TextChunk]:
        """
        Split text into overlapping chunks based on token count

        Args:
            text: The text to chunk
            metadata: Metadata for the chunks
            doc_id: Base document ID

        Returns:
            List of TextChunk objects
        """
        if not text.strip():
            return []

        # Encode the entire text
        tokens = self.encoding.encode(text)
        chunks = []

        start_idx = 0
        chunk_num = 0

        while start_idx < len(tokens):
            # Get chunk tokens
            end_idx = min(start_idx + self.chunk_size, len(tokens))
            chunk_tokens = tokens[start_idx:end_idx]

            # Decode back to text
            chunk_text = self.encoding.decode(chunk_tokens)

            # Create chunk ID
            chunk_id = f"{doc_id}_chunk{chunk_num}"

            # Create chunk object
            chunk = TextChunk(
                id=chunk_id,
                text=chunk_text.strip(),
                metadata=metadata,
                token_count=len(chunk_tokens)
            )
            chunks.append(chunk)

            # Move start index forward by (chunk_size - overlap)
            start_idx += (self.chunk_size - self.chunk_overlap)
            chunk_num += 1

        logger.info(f"Created {len(chunks)} chunks from text (tokens: {len(tokens)})")
        return chunks

    def extract_from_pdf(self, file_path: str) -> List[TextChunk]:
        """
        Extract text from PDF file

        Args:
            file_path: Path to PDF file

        Returns:
            List of TextChunk objects
        """
        logger.info(f"Extracting from PDF: {file_path}")
        chunks = []
        doc_id = self._generate_doc_id(file_path)
        filename = os.path.basename(file_path)

        try:
            # Try pdfplumber first (better text extraction)
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)

                for page_num, page in enumerate(pdf.pages, start=1):
                    text = page.extract_text()

                    if text and text.strip():
                        metadata = ChunkMetadata(
                            filename=filename,
                            doc_type='pdf',
                            page=page_num,
                            total_pages=total_pages,
                            source_path=file_path
                        )

                        page_chunks = self.chunk_text(text, metadata, f"{doc_id}_p{page_num}")
                        chunks.extend(page_chunks)

                logger.info(f"Successfully extracted {len(chunks)} chunks from {total_pages} pages")

        except Exception as e:
            logger.warning(f"pdfplumber failed, trying PyPDF2: {e}")

            try:
                # Fallback to PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    total_pages = len(pdf_reader.pages)

                    for page_num, page in enumerate(pdf_reader.pages, start=1):
                        text = page.extract_text()

                        if text and text.strip():
                            metadata = ChunkMetadata(
                                filename=filename,
                                doc_type='pdf',
                                page=page_num,
                                total_pages=total_pages,
                                source_path=file_path
                            )

                            page_chunks = self.chunk_text(text, metadata, f"{doc_id}_p{page_num}")
                            chunks.extend(page_chunks)

                    logger.info(f"Successfully extracted {len(chunks)} chunks from {total_pages} pages (PyPDF2)")

            except Exception as e2:
                logger.error(f"Failed to extract from PDF {file_path}: {e2}")

        return chunks

    def extract_from_pptx(self, file_path: str) -> List[TextChunk]:
        """
        Extract text from PowerPoint file

        Args:
            file_path: Path to PPTX file

        Returns:
            List of TextChunk objects
        """
        logger.info(f"Extracting from PPTX: {file_path}")
        chunks = []
        doc_id = self._generate_doc_id(file_path)
        filename = os.path.basename(file_path)

        try:
            prs = Presentation(file_path)
            total_slides = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, start=1):
                # Extract text from all shapes in the slide
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                # Also extract notes
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text:
                        slide_text.append(f"\n[Notes: {notes_frame.text}]")

                if slide_text:
                    text = "\n".join(slide_text)

                    metadata = ChunkMetadata(
                        filename=filename,
                        doc_type='pptx',
                        slide=slide_num,
                        total_pages=total_slides,
                        source_path=file_path
                    )

                    slide_chunks = self.chunk_text(text, metadata, f"{doc_id}_s{slide_num}")
                    chunks.extend(slide_chunks)

            logger.info(f"Successfully extracted {len(chunks)} chunks from {total_slides} slides")

        except Exception as e:
            logger.error(f"Failed to extract from PPTX {file_path}: {e}")

        return chunks

    def extract_from_docx(self, file_path: str) -> List[TextChunk]:
        """
        Extract text from Word document

        Args:
            file_path: Path to DOCX file

        Returns:
            List of TextChunk objects
        """
        logger.info(f"Extracting from DOCX: {file_path}")
        chunks = []
        doc_id = self._generate_doc_id(file_path)
        filename = os.path.basename(file_path)

        try:
            doc = Document(file_path)

            # Extract paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        paragraphs.append(row_text)

            if paragraphs:
                text = "\n".join(paragraphs)

                metadata = ChunkMetadata(
                    filename=filename,
                    doc_type='docx',
                    source_path=file_path
                )

                doc_chunks = self.chunk_text(text, metadata, doc_id)
                chunks.extend(doc_chunks)

            logger.info(f"Successfully extracted {len(chunks)} chunks from DOCX")

        except Exception as e:
            logger.error(f"Failed to extract from DOCX {file_path}: {e}")

        return chunks

    def process_file(self, file_path: str) -> List[TextChunk]:
        """
        Process a single file based on its extension

        Args:
            file_path: Path to the file

        Returns:
            List of TextChunk objects
        """
        file_path = str(file_path)
        ext = os.path.splitext(file_path)[1].lower()

        if ext == '.pdf':
            return self.extract_from_pdf(file_path)
        elif ext == '.pptx':
            return self.extract_from_pptx(file_path)
        elif ext == '.docx':
            return self.extract_from_docx(file_path)
        else:
            logger.warning(f"Unsupported file type: {ext} for {file_path}")
            return []

    def process_directory(self, directory: str, recursive: bool = True) -> List[TextChunk]:
        """
        Process all supported files in a directory

        Args:
            directory: Path to directory
            recursive: Whether to process subdirectories

        Returns:
            List of all TextChunk objects
        """
        logger.info(f"Processing directory: {directory} (recursive={recursive})")
        all_chunks = []
        supported_extensions = {'.pdf', '.pptx', '.docx'}

        path = Path(directory)

        if not path.exists():
            logger.error(f"Directory does not exist: {directory}")
            return []

        # Get all files
        if recursive:
            files = [f for f in path.rglob('*') if f.is_file() and f.suffix.lower() in supported_extensions]
        else:
            files = [f for f in path.glob('*') if f.is_file() and f.suffix.lower() in supported_extensions]

        logger.info(f"Found {len(files)} supported files")

        for file_path in files:
            try:
                chunks = self.process_file(str(file_path))
                all_chunks.extend(chunks)
                logger.info(f"Processed {file_path.name}: {len(chunks)} chunks")
            except Exception as e:
                logger.error(f"Error processing {file_path}: {e}")

        return all_chunks

    def _generate_doc_id(self, file_path: str) -> str:
        """Generate a unique document ID from file path"""
        # Use hash of file path for consistent IDs
        hash_obj = hashlib.md5(file_path.encode())
        return hash_obj.hexdigest()[:12]


def main():
    """Main execution function"""

    # Initialize processor
    processor = DocumentProcessor(chunk_size=500, chunk_overlap=50)

    # Define document sources
    document_sources = [
        "/Users/a21/Desktop/Sales Rep Resources 2",
        "/Users/a21/Desktop/RESIDENTIAL_BRAND_GUIDELINES.pdf",
        "/Users/a21/Desktop/Roof-ER Sales Training (1).pptx",
        "/Users/a21/Desktop/Warranty Comparison Prsentation (1).pptx",
        "/Users/a21/Desktop/Sales Trainer Resources/Roof-ER Sales Training.pptx",
        "/Users/a21/Desktop/TA Discussion & Resources/Culture and Commitment.docx.pdf",
        "/Users/a21/Desktop/Sales Trainer Resources",
    ]

    logger.info("=" * 80)
    logger.info("Starting Document Extraction Pipeline")
    logger.info("=" * 80)

    all_chunks = []

    for source in document_sources:
        logger.info(f"\nProcessing source: {source}")

        if os.path.isfile(source):
            # Process single file
            chunks = processor.process_file(source)
            all_chunks.extend(chunks)
        elif os.path.isdir(source):
            # Process directory
            chunks = processor.process_directory(source, recursive=True)
            all_chunks.extend(chunks)
        else:
            logger.warning(f"Source not found: {source}")

    # Convert chunks to dictionaries
    chunks_data = [chunk.to_dict() for chunk in all_chunks]

    # Prepare output
    output_data = {
        "metadata": {
            "total_chunks": len(chunks_data),
            "chunk_size": processor.chunk_size,
            "chunk_overlap": processor.chunk_overlap,
            "sources_processed": len(document_sources),
            "encoding": "cl100k_base"
        },
        "chunks": chunks_data
    }

    # Save to JSON
    output_path = "/Users/a21/Desktop/routellm-chatbot-railway/data/extracted_chunks.json"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(output_data, f, indent=2, ensure_ascii=False)

    logger.info("=" * 80)
    logger.info(f"Extraction Complete!")
    logger.info(f"Total chunks: {len(chunks_data)}")
    logger.info(f"Output saved to: {output_path}")
    logger.info("=" * 80)

    # Print statistics
    doc_types = {}
    for chunk in chunks_data:
        doc_type = chunk['metadata']['doc_type']
        doc_types[doc_type] = doc_types.get(doc_type, 0) + 1

    logger.info("\nChunks by document type:")
    for doc_type, count in doc_types.items():
        logger.info(f"  {doc_type.upper()}: {count} chunks")


if __name__ == "__main__":
    main()
